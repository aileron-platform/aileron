#!/usr/bin/env python3
"""Build a final PPTX deck from generated full-page images."""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "assets"))

import stage_state  # noqa: E402
from stage_state import StageError, resolve_workspace  # noqa: E402

from deck_builder_common import final_page_images, require_deck_ready, resolve_output_path  # noqa: E402


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Build final PPTX deck from generated page PNGs.")
    parser.add_argument("--workspace", type=Path, default=None)
    parser.add_argument("--session-id", required=True)
    parser.add_argument("--output", type=Path, default=None)
    return parser


def _run(args: argparse.Namespace) -> int:
    workspace = resolve_workspace(args.workspace)
    state = require_deck_ready(workspace, args.session_id, "pptx")
    output = resolve_output_path(workspace, args.output, extension="pptx", state=state)
    images = final_page_images(workspace, args.session_id)
    if not images:
        raise stage_state.InvalidTransitionError(
            "no final page images found",
            {"path": str(stage_state.state_path(workspace, args.session_id).parent / "generation")},
        )

    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]
    for image_path in images:
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(image_path), 0, 0, width=presentation.slide_width, height=presentation.slide_height)
    if presentation.slides and len(presentation.slide_layouts) > 0:
        # python-pptx starts with no slides, so no cleanup is normally needed.
        pass
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)
    sys.stdout.write(f"Wrote {output} ({len(images)} slides)\n")
    return 0


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        return _run(args)
    except StageError as exc:
        sys.stderr.write(exc.render() + "\n")
        return exc.exit_code


if __name__ == "__main__":
    raise SystemExit(main())
